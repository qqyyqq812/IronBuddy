"""
IronBuddy 嵌入式系统展示 PPT 构建器
- 只复用信电学院模板的封面页（slide 2）
- 内页自由设计，套用统一视觉基调（深蓝金玻璃拟态）
- 输出：presentation/IronBuddy_嵌入式系统展示_v2.pptx
"""
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path(__file__).parent
ASSETS = ROOT / "assets"
TEMPLATE = Path("/home/qq/docs/PPT/IronBuddy/PPT_ref/信电学院模板.pptx")
OUTPUT = ROOT / "IronBuddy_嵌入式系统展示_v2.pptx"

# 色板（与 build_figures.py 一致）
C_BG = RGBColor(0x0f, 0x17, 0x24)
C_CARD = RGBColor(0x1e, 0x2a, 0x3a)
C_ACCENT = RGBColor(0xd4, 0xa0, 0x4a)
C_PRIMARY = RGBColor(0x4a, 0x90, 0xe2)
C_SUCCESS = RGBColor(0x7e, 0xd3, 0x21)
C_DANGER = RGBColor(0xd0, 0x02, 0x1b)
C_MUTED = RGBColor(0x8b, 0x96, 0xa8)
C_TEXT_LIGHT = RGBColor(0xf5, 0xf7, 0xfa)
C_BORDER = RGBColor(0x2d, 0x3e, 0x54)
C_DARK = RGBColor(0x1a, 0x1a, 0x1a)

FONT_TITLE = "微软雅黑"
FONT_BODY = "微软雅黑"

SW, SH = Inches(13.333), Inches(7.5)


# ──────────────── utilities ────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, line=None, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.08
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    if not shadow:
        # remove shadow
        spPr = shp._element.spPr
        for ef in spPr.findall(qn("a:effectLst")):
            spPr.remove(ef)
        efLst = etree.SubElement(spPr, qn("a:effectLst"))
    return shp


def add_text(slide, x, y, w, h, text, size=14, color=C_TEXT_LIGHT,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_picture(slide, path, x, y, w=None, h=None):
    if w is None and h is None:
        return slide.shapes.add_picture(str(path), x, y)
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def header_bar(slide, title, sub=""):
    """顶部金色标题条 + 装饰线"""
    add_rect(slide, Inches(0), Inches(0), SW, Inches(0.6), fill=C_CARD)
    # 金色左侧色块
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(0.6), fill=C_ACCENT)
    add_text(slide, Inches(0.35), Inches(0.08), Inches(9), Inches(0.45),
             title, size=22, color=C_TEXT_LIGHT, bold=True, font=FONT_TITLE)
    if sub:
        add_text(slide, Inches(9.5), Inches(0.18), Inches(3.7), Inches(0.3),
                 sub, size=11, color=C_MUTED, align=PP_ALIGN.RIGHT)


def footer(slide, page_num, total=16):
    add_text(slide, Inches(0.3), Inches(7.15), Inches(4), Inches(0.3),
             "IronBuddy · 嵌入式系统综合实验", size=9, color=C_MUTED)
    add_text(slide, Inches(9), Inches(7.15), Inches(4), Inches(0.3),
             f"{page_num} / {total}", size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)


def section_tag(slide, tag_text, color=C_ACCENT):
    """章节标签"""
    shp = add_rect(slide, Inches(0.35), Inches(0.75), Inches(1.6), Inches(0.3),
                   fill=color)
    tb = add_text(slide, Inches(0.35), Inches(0.77), Inches(1.6), Inches(0.26),
                  tag_text, size=10, color=C_DARK, bold=True, align=PP_ALIGN.CENTER)


# ──────────────── slide builders ────────────────
def build_cover(prs, template_prs):
    """Slide 1: 使用模板 slide 2 作为封面，填入我们的数据。"""
    # 从 template 中复制 slide 2 为封面
    src_slide = template_prs.slides[1]
    # 简单方法：在空白 layout 上重绘
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, C_BG)

    # 标题
    add_text(slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.5),
             "浙江大学信息与电子工程学院",
             size=16, color=C_MUTED, align=PP_ALIGN.CENTER, font=FONT_TITLE)
    # 分隔线
    line_rect = add_rect(slide, Inches(5.5), Inches(1.4), Inches(2.3), Inches(0.04),
                         fill=C_ACCENT)

    # 主标题
    add_text(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.4),
             "基于端云协同与双模态感知的\n大型智能嵌入式健身教练系统",
             size=36, color=C_TEXT_LIGHT, bold=True,
             align=PP_ALIGN.CENTER, font=FONT_TITLE)

    # 副标题
    add_text(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.5),
             "IronBuddy · 从肌电穿透视觉盲区到多智能体落地",
             size=18, color=C_ACCENT, align=PP_ALIGN.CENTER, bold=True)

    # 封面插图
    cover_img = ASSETS / "G01_cover.png"
    if cover_img.exists():
        add_picture(slide, cover_img, Inches(3.0), Inches(4.0),
                    w=Inches(7.3), h=Inches(2.6))

    # 信息栏
    add_text(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.35),
             "汇报人：[姓名]   ·   指导：[老师]   ·   2026-04-21",
             size=13, color=C_MUTED, align=PP_ALIGN.CENTER)


def build_toc(prs):
    """Slide 2: 目录"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_BG)
    header_bar(slide, "目录")

    items = [
        ("01", "项目设计目标与痛点背景", "CV 盲区 · 设计蓝图"),
        ("02", "硬件生态与无线通信拓扑", "差分 PCB · UDP · /dev/shm"),
        ("03", "核心技术路线：四次关键抉择  ⭐", "视觉 / 融合 / 数据集 / LLM"),
        ("04", "落地实证与系统闭环", "多进程 IPC · MVC · 语音"),
        ("05", "总结与行业展望", "技术栈 · 反思 · 前路"),
    ]
    for i, (num, title, sub) in enumerate(items):
        y = Inches(1.15 + i * 1.05)
        add_rect(slide, Inches(0.6), y, Inches(12.1), Inches(0.85),
                 fill=C_CARD, line=C_BORDER)
        add_rect(slide, Inches(0.6), y, Inches(0.9), Inches(0.85), fill=C_ACCENT)
        add_text(slide, Inches(0.6), y + Inches(0.18), Inches(0.9), Inches(0.5),
                 num, size=26, color=C_DARK, bold=True,
                 align=PP_ALIGN.CENTER, font=FONT_TITLE)
        add_text(slide, Inches(1.8), y + Inches(0.1), Inches(10), Inches(0.45),
                 title, size=18, color=C_TEXT_LIGHT, bold=True, font=FONT_TITLE)
        add_text(slide, Inches(1.8), y + Inches(0.48), Inches(10), Inches(0.35),
                 sub, size=11, color=C_MUTED)

    footer(slide, 2)


def build_chapter_title(prs, num, title, sub, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_BG)
    # 大号编号
    add_text(slide, Inches(1), Inches(2), Inches(5), Inches(2.5),
             num, size=180, color=C_ACCENT, bold=True,
             align=PP_ALIGN.LEFT, font=FONT_TITLE)
    # 分隔竖线
    add_rect(slide, Inches(5.5), Inches(2.3), Inches(0.05), Inches(2.5),
             fill=C_ACCENT)
    add_text(slide, Inches(6), Inches(2.4), Inches(7), Inches(1),
             title, size=32, color=C_TEXT_LIGHT, bold=True, font=FONT_TITLE)
    add_text(slide, Inches(6), Inches(3.5), Inches(7), Inches(1),
             sub, size=15, color=C_MUTED)
    footer(slide, page)


def build_slide_with_image(prs, section, title, bullets, image_name, page,
                           image_side="right", image_ratio=0.55, caption=""):
    """通用版式：标题条 + 左正文 + 右大图（或反过来）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_BG)
    header_bar(slide, title, section)

    text_x = Inches(0.5)
    text_w = Inches(5.3)
    img_x = Inches(6.0)
    img_w = Inches(7.0)
    if image_side == "left":
        img_x = Inches(0.3)
        text_x = Inches(7.3)
        text_w = Inches(5.7)
        img_w = Inches(6.8)

    # 正文条目
    y = Inches(1.05)
    for b in bullets:
        if isinstance(b, tuple):
            head, body = b
            # 卡片
            add_rect(slide, text_x, y, text_w, Inches(1.15),
                     fill=C_CARD, line=C_BORDER)
            add_text(slide, text_x + Inches(0.15), y + Inches(0.08),
                     text_w - Inches(0.3), Inches(0.35),
                     head, size=13, color=C_ACCENT, bold=True, font=FONT_TITLE)
            add_text(slide, text_x + Inches(0.15), y + Inches(0.45),
                     text_w - Inches(0.3), Inches(0.65),
                     body, size=11, color=C_TEXT_LIGHT, line_spacing=1.25)
            y += Inches(1.3)
        else:
            add_text(slide, text_x, y, text_w, Inches(0.5),
                     "• " + b, size=13, color=C_TEXT_LIGHT)
            y += Inches(0.55)

    # 图片
    img_path = ASSETS / image_name
    if img_path.exists():
        img_h = Inches(5.4)
        add_picture(slide, img_path, img_x, Inches(1.3), w=img_w, h=img_h)
        if caption:
            add_text(slide, img_x, Inches(6.75), img_w, Inches(0.3),
                     caption, size=9, color=C_MUTED,
                     align=PP_ALIGN.CENTER, font=FONT_BODY)

    footer(slide, page)


def build_slide_full_image(prs, section, title, image_name, caption, page):
    """满屏图版式"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_BG)
    header_bar(slide, title, section)
    img_path = ASSETS / image_name
    if img_path.exists():
        add_picture(slide, img_path, Inches(0.5), Inches(1.0),
                    w=Inches(12.3), h=Inches(5.8))
    add_text(slide, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.35),
             caption, size=11, color=C_MUTED, align=PP_ALIGN.CENTER)
    footer(slide, page)


def build_slide_text_heavy(prs, section, title, sections, page):
    """纯文字版式（用于技术对比表）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_BG)
    header_bar(slide, title, section)

    y = Inches(1.05)
    for sec in sections:
        if isinstance(sec, dict):
            head = sec.get("head", "")
            body = sec.get("body", [])
            color = sec.get("color", C_ACCENT)
            add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.4),
                     fill=color)
            add_text(slide, Inches(0.7), y + Inches(0.05), Inches(12),
                     Inches(0.3), head, size=13, color=C_DARK, bold=True)
            y += Inches(0.5)
            for line in body:
                add_text(slide, Inches(0.8), y, Inches(12), Inches(0.35),
                         "• " + line, size=12, color=C_TEXT_LIGHT)
                y += Inches(0.4)
            y += Inches(0.1)
    footer(slide, page)


# ──────────────── build ────────────────
def build():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    template_prs = Presentation(str(TEMPLATE))

    # ============ Slide 1: 封面 ============
    build_cover(prs, template_prs)

    # ============ Slide 2: 目录 ============
    build_toc(prs)

    # ============ Slide 3: 行业痛点 ============
    build_slide_with_image(
        prs, section="第1章 · 背景",
        title="行业痛点：纯视觉感知的深层盲区",
        bullets=[
            ("骨架一致 ≠ 发力正确",
             "市面健身镜仅做 CV 骨架追踪，只能看「形」不能看「劲」"),
            ("代偿借力是 CV 看不见的伤",
             "同一下蹲轨迹：股四卸力 / 腓肠代偿 → 韧带隐性慢损"),
            ("缺口：生物电维度补足视觉",
             "sEMG 穿透皮肤读取微伏级激活，是唯一"),
        ],
        image_name="G03_cv_blind.png", page=3,
        caption="左：骨架合规  |  右：肌电暴露深层代偿")

    # ============ Slide 4: 设计蓝图 ============
    build_slide_with_image(
        prs, section="第1章 · 背景",
        title="IronBuddy 系统设计蓝图",
        bullets=[
            ("感知双模", "视觉骨架（宏观）+ sEMG 肌电（深层）"),
            ("算力双端", "板端 NPU 零延迟 + 云端 RTX 5090 高精度"),
            ("认知双轨", "DeepSeek 短视实时 + OpenClaw 常驻长视"),
            ("交互无触", "玻璃拟态只读 UI + 百度 AipSpeech 语音全控"),
        ],
        image_name="G04_architecture.png", page=4,
        caption="五层架构：硬件 / 通信 / 感知 / 认知 / 交互")

    # ============ Slide 5: 硬件生态 ============
    build_slide_with_image(
        prs, section="第2章 · 硬件",
        title="硬件生态：从杂波崩溃到穿戴解耦",
        bullets=[
            ("定制差分 PCB",
             "双路 sEMG + ESP32 + 独立锂电 · 单点地隔离 · 底噪 ≤ 10 μV"),
            ("穿戴式腰包",
             "主控挂腰 · 医用硅胶线 · 摒弃杜邦线束缚"),
            ("量化指标",
             "360° 佩戴自由 · 电极 ≤ 2 min 贴附 · 实验室级信噪比"),
        ],
        image_name="S01.png", page=5,
        caption="S01：PCB 实物（占位图 · 明早替换棚拍）",
        image_side="right")

    # ============ Slide 6: 通信拓扑 ============
    build_slide_with_image(
        prs, section="第2章 · 硬件",
        title="无线通信拓扑：UDP 透传 + /dev/shm 心跳",
        bullets=[
            ("对外 WiFi UDP",
             "ESP32 → ~1 kHz ASCII 封包 · 双网卡隔离"),
            ("对内 /dev/shm",
             "5 进程 × 20+ 信号文件 · atomic rename 保证一致性"),
            ("绕 GIL 锁",
             "视觉进程内嵌 MJPEG:8080 · HDMI 后 Flask CPU 从 30% → <5%"),
        ],
        image_name="G05_comm_topology.png", page=6,
        caption="左：WiFi UDP 对外  |  右：/dev/shm 对内辐射")

    # ============ Slide 7: 心路①视觉 ============
    build_slide_with_image(
        prs, section="第3章 · 心路 ① 视觉",
        title="本地 NPU vs 云端 GPU —— 不是二选一，是热切换",
        bullets=[
            ("本地 NPU（默认）",
             "YOLOv5-Pose RKNN uint8 · conf=0.08 · ~107 ms/帧 · 零网络"),
            ("云端 GPU（可切）",
             "RTMPose-m ONNX · RTX 5090 · ~30 ms + RTT · 高精度"),
            ("热切换实现",
             "写 /dev/shm/vision_mode.json 即切 · 云端超时自动降级"),
            ("心路教训",
             "量化模型置信度最高 0.2，用 0.5 阈值完全看不到人"),
        ],
        image_name="G06_vision_dual.png", page=7,
        caption="双引擎 · 信号文件驱动")

    # ============ Slide 8: 心路②融合 ============
    build_slide_with_image(
        prs, section="第3章 · 心路 ② 融合",
        title="Mid-Fusion：端到端 OOM → 1488 参数通吃",
        bullets=[
            ("❌ 端到端双分支",
             "原始张量直接塞注意力网络 · 模态不对齐 · 4GB 几秒 OOM"),
            ("✅ 中继融合",
             "骨架→角度/相位 · EMG→Biquad+FFT+RMS · 7D 标量"),
            ("微型 GRU(hidden=16)",
             "1488 参数 · <1 MB · 3 头输出 · CPU 微秒级"),
            ("心路教训",
             "边缘算力不是无底洞 —— 这是整个项目最大的课"),
        ],
        image_name="G07_mid_fusion.png", page=8,
        caption="领域知识先降维，再喂微型 GRU")

    # ============ Slide 9: 心路③数据集 ============
    build_slide_with_image(
        prs, section="第3章 · 心路 ③ 数据集",
        title="数据集抉择：3 天的时间线",
        bullets=[
            ("Ninapro / Camargo ✗",
             "抓握 / 步态 · 动作不对口"),
            ("FLEX NeurIPS25 ✗",
             "License 审批 24-72h · 明日等不到"),
            ("MIA ICCV23 ✓",
             "964 clip 下肢 → 深蹲 val acc 94.4%"),
            ("本地自采 + 10× augment ✓",
             "弯举兜底 · 动态 MVC · val acc 94→100%"),
        ],
        image_name="G08_dataset_tree.png", page=9,
        caption="本地化不是妥协，与 ESP32 硬件域完美对齐")

    # ============ Slide 10: 心路④双轨LLM ============
    build_slide_with_image(
        prs, section="第3章 · 心路 ④ LLM",
        title="双轨双核 LLM：实时 × 常驻",
        bullets=[
            ("DeepSeek 前端实时",
             "语音 / 按钮触发 · 短视无历史 · SSE 流式 · ≤ 2 s"),
            ("OpenClaw 后端常驻",
             "cron 09/20/23 · 全量 14 日上下文 · 偏好学习"),
            ("SQLite 8 表共享",
             "rep_events / llm_log / preference_history / voice_sessions..."),
            ("心路教训",
             "一个 LLM 做不了实时陪练 + 长期记忆两件事"),
        ],
        image_name="G09_dual_llm.png", page=10,
        caption="延迟敏感 vs 无延迟约束 · 分工即解耦")

    # ============ Slide 11: 实证① 5进程 IPC ============
    build_slide_with_image(
        prs, section="第4章 · 实证",
        title="5 进程 + /dev/shm 原子 IPC",
        bullets=[
            ("5 进程独立部署",
             "vision / streamer / fsm / emg / voice · 崩一个不影响其他"),
            ("20+ 信号文件",
             "pose_data / fsm_state / muscle_activation / chat_input ..."),
            ("原子写入",
             "tmp → rename POSIX 语义 · 消费者永不读半写数据"),
            ("效果",
             "HDMI 激活后 Flask CPU 从 30% 降到 < 5%"),
        ],
        image_name="G10_ipc_mesh.png", page=11,
        caption="辐射网络 · 每条边是一组信号文件")

    # ============ Slide 12: 实证② MVC 校准 ============
    build_slide_with_image(
        prs, section="第4章 · 实证",
        title="MVC 动态校准 + 硬件域对齐",
        bullets=[
            ("个体归一化：3.5 s MVC 峰值",
             "用户最大发力 → 取峰值 RMS → 写 mvc_values.json"),
            ("硬件域对齐：α·x+β",
             "ESP32 vs MIA Delsys 域差异 → 本地 α=2.12 β=-21.4"),
            ("不动 MIA 权重",
             "val acc 94.4% 不退化 · 输入侧线性滤镜"),
            ("心路教训",
             "硬编码 400 是对人体异质性的侮辱"),
        ],
        image_name="G11_alpha_beta.png", page=12,
        caption="p05/p95 对齐 · 线性映射零损失")

    # ============ Slide 13: 实证③ 语音 ============
    build_slide_with_image(
        prs, section="第4章 · 实证",
        title="语音交互：Vosk ABI → edge-tts 网络 → 百度 AipSpeech",
        bullets=[
            ("❌ Vosk 本地离线",
             "Debian 10 ARM64 · glibc ABI 不兼容 · libvosk.so 符号缺失"),
            ("❌ edge-tts 微软云",
             "板端网络抖动时完全不可用"),
            ("✅ 百度 AipSpeech",
             "稳定 TTS/STT · 延迟 < 500 ms · 免费额度够用"),
            ("关键特性",
             "自适应 VAD · ALSA 日志屏蔽 · 4 条语音命令"),
        ],
        image_name="S07.png", page=13,
        caption="S07：PWA 前端完整截图（占位图 · 明早替换）")

    # ============ Slide 14: 总结·技术栈 ============
    build_slide_with_image(
        prs, section="第5章 · 总结",
        title="全栈技术闭环 + 知识盲区反思",
        bullets=[
            ("✓ 技术栈 5 层闭环",
             "硬件 / 信号 / AI / 系统 / 应用 —— 从 C 到 PWA"),
            ("盲区 ①：算力规划缺位",
             "原型期端到端 OOM —— 最大教训"),
            ("盲区 ②：通信队列薄弱",
             "Python GIL 卡顿 → 才定位到 /dev/shm 非阻塞心跳"),
        ],
        image_name="G12_pyramid.png", page=14,
        caption="每一层都踩过坑，每一层都有文档追溯")

    # ============ Slide 15: 行业展望 ============
    build_slide_with_image(
        prs, section="第5章 · 展望",
        title="行业展望：双模态感知的 3 条前路",
        bullets=[
            ("🩺 柔性印刷电极阵列",
             "2 通道 → N 通道 · 全身肌群网络 · 术后复健无创量化"),
            ("📊 开源错误动作数据集",
             "业界当前空白 · 本项目 augment 沙盒即为雏形"),
            ("🧠 边缘端侧大模型",
             "期待 7B 量化落地 NPU · 腰包即私人教练"),
        ],
        image_name="G13_outlook.png", page=15,
        caption="从健身工具 → 医疗复健 → 基础研究")

    # ============ Slide 16: Q&A ============
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_BG)
    add_picture(slide, ASSETS / "G14_thanks.png", Inches(0), Inches(0),
                w=SW, h=SH)

    # 保存
    prs.save(str(OUTPUT))
    print(f"✓ 输出：{OUTPUT}")
    print(f"  文件大小：{OUTPUT.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    build()
